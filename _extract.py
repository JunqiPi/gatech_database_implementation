# -*- coding: utf-8 -*-
"""Extract all text from PPTX and PDF files into Markdown, verbatim.

paper1.pdf is handled separately by _reflow_paper1.py (two-column reflow),
so it is skipped here.
"""
import os, glob
from pptx import Presentation
import fitz  # PyMuPDF

BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "extracted")
os.makedirs(OUT, exist_ok=True)

# PowerPoint stores Symbol/Wingdings-font glyphs in the Unicode Private Use
# Area (U+F000 + char code). Map the ones actually used in these decks back
# to their real Unicode math symbols so relational-algebra notation survives.
SYMBOL_MAP = {
    chr(0xF02A): "×",  # x  multiplication  (n x (i-1))
    chr(0xF072): "ρ",  # rho   -> rename
    chr(0xF073): "σ",  # sigma -> select
    chr(0xF0A3): "≤",  # <=
    chr(0xF0AC): "←",  # <-   assignment
    chr(0xF0B3): "≥",  # >=
    chr(0xF0B9): "≠",  # !=
    chr(0xF0C7): "∩",  # intersection
    chr(0xF0C8): "∪",  # union
    chr(0xF0CD): "⊆",  # subset-or-equal
    chr(0xF0D5): "Π",  # Pi    -> project
    chr(0xF0D8): "¬",  # not
    chr(0xF0D9): "∧",  # and
    chr(0xF0DA): "∨",  # or
    chr(0xF0E8): "→",  # ->   (Wingdings arrow: "results in")
}


def fix_symbols(text):
    if not text:
        return text
    for k, v in SYMBOL_MAP.items():
        if k in text:
            text = text.replace(k, v)
    return text


def iter_shapes(shapes):
    """Recursively yield shapes, descending into groups."""
    for shp in shapes:
        if shp.shape_type == 6:  # GROUP
            yield from iter_shapes(shp.shapes)
        else:
            yield shp


def table_to_md(table):
    rows = []
    for r in table.rows:
        cells = [fix_symbols(c.text).replace("\n", " ").strip() for c in r.cells]
        rows.append(cells)
    if not rows:
        return ""
    out = []
    header = rows[0]
    out.append("| " + " | ".join(header) + " |")
    out.append("| " + " | ".join(["---"] * len(header)) + " |")
    for r in rows[1:]:
        r = (r + [""] * len(header))[:len(header)]
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


def extract_pptx(path):
    prs = Presentation(path)
    name = os.path.splitext(os.path.basename(path))[0]
    lines = [f"# {name}", ""]
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        lines.append("")
        body_blocks = []
        for shp in iter_shapes(slide.shapes):
            if shp.has_table:
                md = table_to_md(shp.table)
                if md:
                    body_blocks.append(md)
                continue
            if shp.has_text_frame:
                txt = fix_symbols(shp.text_frame.text)
                if txt and txt.strip():
                    body_blocks.append(txt.rstrip())
        notes = ""
        if slide.has_notes_slide:
            n = fix_symbols(slide.notes_slide.notes_text_frame.text)
            if n and n.strip():
                notes = n.rstrip()
        for b in body_blocks:
            lines.append(b)
            lines.append("")
        if notes:
            lines.append("> **[备注 Notes]**")
            for nl in notes.splitlines():
                lines.append(f"> {nl}")
            lines.append("")
        lines.append("")
    out_path = os.path.join(OUT, name + ".md")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines).rstrip() + "\n")
    return out_path, len(prs.slides)


def extract_pdf(path):
    doc = fitz.open(path)
    name = os.path.splitext(os.path.basename(path))[0]
    lines = [f"# {name}", ""]
    for i, page in enumerate(doc, 1):
        lines.append(f"## Page {i}")
        lines.append("")
        txt = fix_symbols(page.get_text("text"))
        lines.append(txt.rstrip())
        lines.append("")
    out_path = os.path.join(OUT, name + ".md")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines).rstrip() + "\n")
    return out_path, doc.page_count


if __name__ == "__main__":
    print("=== PPTX ===")
    for p in sorted(glob.glob(os.path.join(BASE, "*.pptx"))):
        op, n = extract_pptx(p)
        print(f"{os.path.basename(p)} -> {os.path.basename(op)} ({n} slides)")
    print("=== PDF ===")
    for p in sorted(glob.glob(os.path.join(BASE, "*.pdf"))):
        if os.path.basename(p).lower() == "paper1.pdf":
            print("paper1.pdf -> handled by _reflow_paper1.py (skipped)")
            continue
        op, n = extract_pdf(p)
        print(f"{os.path.basename(p)} -> {os.path.basename(op)} ({n} pages)")
    print("Done. Output in:", OUT)
